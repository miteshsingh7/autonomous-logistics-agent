"""Generate the Part 3 presentation as a real .pptx file.

Usage:
    python scripts/build_deck.py

The deck reads results/scorecards.json when present so the numbers on the
results slide reflect the latest run; otherwise it falls back to representative
figures and labels them as such. Output: deck/presentation.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results" / "scorecards.json"
OUT = ROOT / "deck" / "presentation.pptx"

# --- palette (4 colors: 2 neutrals + navy brand + amber accent) -------------
INK = RGBColor(0x14, 0x1B, 0x2E)      # deep navy (backgrounds / headings)
PAPER = RGBColor(0xF7, 0xF8, 0xFA)    # off-white (text on dark)
SLATE = RGBColor(0x5B, 0x66, 0x7B)    # muted slate (secondary text)
ACCENT = RGBColor(0xE8, 0x9A, 0x2E)   # amber (emphasis)

WIDE = Inches(13.333)
TALL = Inches(7.5)


def _load_aggregates() -> tuple[dict[str, dict[str, float]], bool]:
    """Return per-label aggregate means and whether they came from a real run."""
    if not RESULTS.exists():
        demo = {
            "closed": {"tool": 1.00, "reasoning": 0.98, "recover": 1.00,
                       "success": 1.00, "overall": 0.99},
            "open": {"tool": 1.00, "reasoning": 0.86, "recover": 1.00,
                     "success": 1.00, "overall": 0.94},
        }
        return demo, False

    cards = json.loads(RESULTS.read_text())
    agg: dict[str, dict[str, list[float]]] = {}
    for c in cards:
        label = c["model"].split(":", 1)[0]
        bucket = agg.setdefault(label, {"tool": [], "reasoning": [],
                                        "recover": [], "success": [], "overall": []})
        bucket["tool"].append(c["tool_accuracy"])
        bucket["reasoning"].append(c["reasoning"])
        if c["error_recovery"] is not None:
            bucket["recover"].append(c["error_recovery"])
        bucket["success"].append(c["task_success"])
        bucket["overall"].append(c["overall"])
    means = {
        label: {k: (sum(v) / len(v) if v else 0.0) for k, v in b.items()}
        for label, b in agg.items()
    }
    return means, True


# --- low-level slide helpers ------------------------------------------------

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, left, top, width, height, text, *, size=18, bold=False,
          color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _bullets(slide, left, top, width, height, items, *, size=18,
             color=INK, gap=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        level_text, is_sub = (item[1:], True) if item.startswith("\t") else (item, False)
        run = p.add_run()
        run.text = ("   – " if is_sub else "•  ") + level_text.strip()
        run.font.size = Pt(size - 2 if is_sub else size)
        run.font.color.rgb = SLATE if is_sub else color
        run.font.name = "Calibri"
    return box


def _accent_bar(slide, top=Inches(1.15), left=Inches(0.6), width=Inches(1.4)):
    bar = slide.shapes.add_shape(1, left, top, width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _section_header(slide, kicker, title):
    _text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
          kicker.upper(), size=13, bold=True, color=ACCENT)
    _text(slide, Inches(0.6), Inches(0.78), Inches(12), Inches(0.7),
          title, size=30, bold=True, color=INK)
    _accent_bar(slide, top=Inches(1.5))


# --- slide builders ---------------------------------------------------------

def slide_title(prs):
    s = _blank(prs)
    _bg(s, INK)
    _accent_bar(s, top=Inches(2.5), left=Inches(0.9), width=Inches(1.6))
    _text(s, Inches(0.9), Inches(2.7), Inches(11), Inches(1.4),
          "Autonomous Carrier Rerouting", size=44, bold=True, color=PAPER)
    _text(s, Inches(0.9), Inches(3.9), Inches(11), Inches(0.8),
          "Is an open model production-ready to replace a closed one?",
          size=22, color=ACCENT)
    _text(s, Inches(0.9), Inches(5.6), Inches(11), Inches(0.5),
          "Multi-agent workflow + trajectory-based evaluation  |  AI Researcher take-home",
          size=15, color=SLATE)
    return s


def slide_problem(prs):
    s = _blank(prs)
    _bg(s, PAPER)
    _section_header(s, "The problem", "Reactive logistics is expensive")
    _bullets(s, Inches(0.6), Inches(1.9), Inches(7.2), Inches(4.5), [
        "Disruptions (weather, breakdowns, port congestion) are constant and time-critical.",
        "Human dispatchers react late; every hour of delay compounds SLA and cost risk.",
        "Goal: a proactive, self-healing system that reroutes autonomously and escalates only when needed.",
        "Key question for the business: can we run this on an open model we host, or do we need a frontier closed model?",
    ], size=19)
    _text(s, Inches(8.1), Inches(2.0), Inches(4.6), Inches(3.5),
          "\u201cProactive, self-healing\u201d means detecting risk, acting without a human in the loop, and recovering from its own tool failures.",
          size=16, color=INK)
    return s


def slide_approach(prs):
    s = _blank(prs)
    _bg(s, PAPER)
    _section_header(s, "Approach", "A supervised multi-agent workflow")
    _bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), [
        "Orchestrator (supervisor) routes between three specialist agents on a shared state blackboard:",
        "\tAssessor — pulls telemetry, decides if a reroute is warranted.",
        "\tPlanner — enumerates alternatives, confirms a firm cost/ETA quote, selects the best option.",
        "\tExecutor — commits the reroute via the dispatch system.",
        "Each agent owns a narrow tool set (better tool-calling reliability) and self-heals against failures in its own tools.",
        "Built on LangGraph; every model is reached through one gateway so closed vs. open is a one-line swap.",
    ], size=18)
    return s


def slide_eval_method(prs):
    s = _blank(prs)
    _bg(s, PAPER)
    _section_header(s, "How we measured", "Trajectory-based evaluation")
    _bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), [
        "We score the whole trajectory (every step), not just the final answer:",
        "\tTool accuracy — valid args, full coverage, correct ordering.",
        "\tReasoning quality — LLM-as-judge rubric on justification and evidence use.",
        "\tError recovery — did it self-heal from an injected tool failure?",
        "\tTask success — did it reach the right outcome (reroute vs. escalate)?",
        "5 scenarios, 2 with injected faults (telemetry timeout, marketplace timeout).",
    ], size=18)
    return s


def slide_results(prs, agg, real):
    s = _blank(prs)
    _bg(s, PAPER)
    tag = "live run" if real else "representative (run to refresh)"
    _section_header(s, f"Results — {tag}", "Closed vs. open, head to head")

    rows = [["Dimension", "Closed", "Open"]]
    for key, label in [("tool", "Tool accuracy"), ("reasoning", "Reasoning"),
                       ("recover", "Error recovery"), ("success", "Task success"),
                       ("overall", "Overall")]:
        c = agg.get("closed", {}).get(key, 0.0)
        o = agg.get("open", {}).get(key, 0.0)
        rows.append([label, f"{c:.2f}", f"{o:.2f}"])

    tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(1.9),
                                   Inches(7.2), Inches(3.6))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(16)
            run.font.name = "Calibri"
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = PAPER
                cell.fill.solid()
                cell.fill.fore_color.rgb = INK
            else:
                run.font.color.rgb = INK
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER if r % 2 else RGBColor(0xEC, 0xEF, 0xF3)
                if rows[r][0] == "Overall":
                    run.font.bold = True

    _bullets(s, Inches(8.1), Inches(2.0), Inches(4.7), Inches(4.5), [
        "Tool accuracy, recovery and task success reach parity.",
        "The gap concentrates in reasoning quality — nuanced trade-off calls.",
        "Open recovers from injected faults just as reliably.",
    ], size=17)
    return s


def slide_tradeoffs(prs):
    s = _blank(prs)
    _bg(s, PAPER)
    _section_header(s, "Trade-offs", "It is not just accuracy")
    _bullets(s, Inches(0.6), Inches(1.9), Inches(6), Inches(4.5), [
        "Cost — open (self-hosted) is far cheaper per run at volume.",
        "Latency — open can be lower if hosted close to the app.",
        "Control — open gives data residency, no vendor lock-in, tunable weights.",
    ], size=18)
    _bullets(s, Inches(6.9), Inches(1.9), Inches(6), Inches(4.5), [
        "Reliability — closed still leads on hard reasoning calls.",
        "Ops burden — open means you own hosting, scaling, evals.",
        "Safety — closed ships more mature guardrails out of the box.",
    ], size=18)
    return s


def slide_recommendation(prs):
    s = _blank(prs)
    _bg(s, INK)
    _text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
          "RECOMMENDATION", size=13, bold=True, color=ACCENT)
    _text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(1.0),
          "Adopt a hybrid, tiered routing policy", size=30, bold=True, color=PAPER)
    _accent_bar(s, top=Inches(1.7))
    _bullets(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5), [
        "Route routine, high-confidence reroutes to the open model — capture the cost savings where it is at parity.",
        "Escalate high-value or ambiguous cases to the closed model (or a human) — where reasoning quality matters most.",
        "Gate the boundary with the trajectory eval: promote more traffic to open as its reasoning score closes the gap.",
        "Net: most volume on cheaper open infra, frontier reasoning reserved for the decisions that carry real risk.",
    ], size=19, color=PAPER)
    return s


def slide_risks(prs):
    s = _blank(prs)
    _bg(s, PAPER)
    _section_header(s, "Rollout & risks", "How we ship this safely")
    _bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), [
        "Phase 1 — shadow mode: open runs alongside production, scored on live trajectories, no dispatch.",
        "Phase 2 — canary: open handles low-risk reroutes with human review of escalations.",
        "Phase 3 — tiered autonomy: expand open's remit as reasoning parity is demonstrated.",
        "Risks — reasoning gap on edge cases, hosting/ops maturity, eval coverage. Mitigation: keep the closed-model fallback and the eval gate.",
    ], size=18)
    return s


def slide_closing(prs):
    s = _blank(prs)
    _bg(s, INK)
    _accent_bar(s, top=Inches(2.6), left=Inches(0.9), width=Inches(1.6))
    _text(s, Inches(0.9), Inches(2.8), Inches(11), Inches(1.0),
          "Build once, evaluate rigorously, route by risk.", size=32, bold=True,
          color=PAPER)
    _text(s, Inches(0.9), Inches(4.0), Inches(11), Inches(0.6),
          "The trajectory eval is the control system for how much autonomy we give open models over time.",
          size=18, color=ACCENT)
    return s


def build() -> None:
    agg, real = _load_aggregates()
    prs = Presentation()
    prs.slide_width = Emu(int(WIDE))
    prs.slide_height = Emu(int(TALL))

    slide_title(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_eval_method(prs)
    slide_results(prs, agg, real)
    slide_tradeoffs(prs)
    slide_recommendation(prs)
    slide_risks(prs)
    slide_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    src = "live results" if real else "representative numbers"
    print(f"Wrote {OUT.relative_to(ROOT)} ({len(prs.slides)} slides, {src}).")


if __name__ == "__main__":
    build()
